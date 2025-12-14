import os
import time
import streamlit as st

from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

import google.generativeai as genai

from langchain_community.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_google_genai import (
    GoogleGenerativeAIEmbeddings,
    ChatGoogleGenerativeAI
)
from langchain.prompts import PromptTemplate
from langchain.chains.question_answering import load_qa_chain

# -------------------------------------------------------------------
# ENV SETUP
# -------------------------------------------------------------------
load_dotenv()

API_KEY = os.getenv("GOOGLE_API_KEY")
if not API_KEY:
    st.error("GOOGLE_API_KEY is not set")
    st.stop()

genai.configure(api_key=API_KEY)

# -------------------------------------------------------------------
# FILE TEXT EXTRACTION
# -------------------------------------------------------------------
def extract_text_from_pdf(file):
    text = ""
    reader = PdfReader(file)
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text
    return text


def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_text_from_pptx(file):
    text = ""
    presentation = Presentation(file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text += para.text + "\n"
    return text


def extract_text_from_txt(file):
    return file.read().decode("utf-8")


def extract_text_from_file(file):
    if file.name.endswith(".pdf"):
        return extract_text_from_pdf(file)
    if file.name.endswith(".docx"):
        return extract_text_from_docx(file)
    if file.name.endswith(".pptx"):
        return extract_text_from_pptx(file)
    if file.name.endswith(".txt"):
        return extract_text_from_txt(file)
    return ""

# -------------------------------------------------------------------
# TEXT CHUNKING & VECTOR STORE
# -------------------------------------------------------------------
def split_text(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=10000,
        chunk_overlap=1000
    )
    return splitter.split_text(text)


@st.cache_resource
def load_embeddings():
    return GoogleGenerativeAIEmbeddings(model="models/embedding-001")

# -------------------------------------------------------------------
# QA CHAIN
# -------------------------------------------------------------------
def get_qa_chain():
    prompt = PromptTemplate(
        template="""
Answer the question using ONLY the provided context.
If the answer is not present, say "I don't know".

Context:
{context}

Question:
{question}

Answer:
""",
        input_variables=["context", "question"]
    )

    model = ChatGoogleGenerativeAI(
        model="gemini-pro",
        temperature=0.3
    )

    return load_qa_chain(
        llm=model,
        chain_type="stuff",
        prompt=prompt
    )

# -------------------------------------------------------------------
# USER QUERY HANDLER
# -------------------------------------------------------------------
def handle_user_question(question):
    if "vector_store" not in st.session_state:
        return "Please upload and process files first."

    if "qa_chain" not in st.session_state:
        st.session_state.qa_chain = get_qa_chain()

    docs = st.session_state.vector_store.similarity_search(question, k=4)

    result = st.session_state.qa_chain.invoke(
        {"input_documents": docs, "question": question},
        return_only_outputs=True
    )

    return result["output_text"]

# -------------------------------------------------------------------
# STREAMLIT UI
# -------------------------------------------------------------------
def main():
    st.set_page_config(page_title="Chat With Files", layout="wide")
    st.header("Chat with Files")

    if "conversation" not in st.session_state:
        st.session_state.conversation = []

    with st.sidebar:
        st.title("Menu")

        uploaded_files = st.file_uploader(
            "Upload PDF, DOCX, PPTX, or TXT files",
            accept_multiple_files=True,
            type=["pdf", "docx", "pptx", "txt"]
        )

        if st.button("Submit & Process"):
            if not uploaded_files:
                st.error("Please upload at least one file.")
            else:
                with st.spinner("Processing files..."):
                    raw_text = ""

                    for file in uploaded_files:
                        if file.size > 10 * 1024 * 1024:
                            st.error(f"{file.name} is too large (max 10MB).")
                            return
                        raw_text += extract_text_from_file(file)

                    if not raw_text.strip():
                        st.error("No readable text found.")
                        return

                    chunks = split_text(raw_text)
                    embeddings = load_embeddings()

                    # IMPORTANT: keep FAISS in session (no disk writes)
                    st.session_state.vector_store = FAISS.from_texts(
                        chunks, embedding=embeddings
                    )

                    st.success("Files processed. You can now ask questions.")

    # Chat history
    for q, a in st.session_state.conversation:
        st.markdown(f"**You:** {q}")
        st.markdown(f"**Model:** {a}")

    user_question = st.text_input("Ask a question")

    if st.button("Submit Query"):
        if user_question:
            answer = handle_user_question(user_question)
            st.session_state.conversation.append((user_question, answer))
            st.rerun()

# -------------------------------------------------------------------
if __name__ == "__main__":
    main()
